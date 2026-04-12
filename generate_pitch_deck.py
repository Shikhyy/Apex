#!/usr/bin/env python3
"""
Generate APEX Pitch Deck as PowerPoint (PPTX)
Converts the JSX pitch deck to professional PowerPoint format
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# Color scheme from JSX
COLORS = {
    "orange": RGBColor(0xD5, 0x3E, 0x0F),
    "dark": RGBColor(0x9B, 0x0F, 0x06),
    "deepRed": RGBColor(0x5E, 0x00, 0x06),
    "cream": RGBColor(0xEE, 0xD9, 0xB9),
    "black": RGBColor(0x0D, 0x0A, 0x09),
    "charcoal": RGBColor(0x1A, 0x0F, 0x0D),
    "dimRed": RGBColor(0x2A, 0x0A, 0x08),
    "green": RGBColor(0x34, 0xD3, 0x99),
}

def add_background(slide, color):
    """Add solid background color to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_with_highlight(slide, title_text, highlight_word=None, y_pos=None):
    """Add title with optional highlighted text"""
    if y_pos is None:
        y_pos = Inches(1.5)
    
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(2)
    
    text_box = slide.shapes.add_textbox(left, y_pos, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # Split by <span> tags if highlight_word exists
    if highlight_word:
        parts = title_text.split(highlight_word)
        for i, part in enumerate(parts):
            if i > 0:
                p = text_frame.add_paragraph() if i > 1 else text_frame.paragraphs[0]
                run = p.add_run() if i > 1 else p.runs[0] if p.runs else p.add_run()
                run.text = highlight_word
                run.font.size = Pt(60)
                run.font.bold = True
                run.font.color.rgb = COLORS["orange"]
                
                if i < len(parts) - 1:
                    if len(p.runs) == 1:
                        p.text = part
                    else:
                        new_p = text_frame.add_paragraph()
                        new_p.text = parts[i + 1] if i + 1 < len(parts) else ""
            else:
                if i == 0:
                    p = text_frame.paragraphs[0]
                    p.text = part
                    p.font.size = Pt(60)
                    p.font.bold = True
                    p.font.color.rgb = COLORS["cream"]
    else:
        p = text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = COLORS["cream"]

def add_subtitle(slide, text, y_pos=None):
    """Add subtitle text"""
    if y_pos is None:
        y_pos = Inches(4)
    
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1.5)
    
    text_box = slide.shapes.add_textbox(left, y_pos, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS["dark"]
    p.line_spacing = 1.5

def add_tag(slide, text, x_pos=0.5, y_pos=0.3):
    """Add tag element"""
    left = Inches(x_pos)
    top = Inches(y_pos)
    width = Inches(2)
    height = Inches(0.3)
    
    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["dimRed"]
    shape.line.color.rgb = COLORS["dark"]
    shape.line.width = Pt(1)
    
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["orange"]
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_bullet_list(slide, items, x_pos=0.5, y_pos=2, width=9):
    """Add bullet list"""
    left = Inches(x_pos)
    top = Inches(y_pos)
    list_width = Inches(width)
    list_height = Inches(len(items) * 0.5 + 0.5)
    
    text_box = slide.shapes.add_textbox(left, top, list_width, list_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = "◈ " + item
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["cream"]
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(4)

def add_divider(slide, x_pos=0.5, y_pos=2):
    """Add orange divider line"""
    left = Inches(x_pos)
    top = Inches(y_pos)
    width = Inches(0.5)
    height = Inches(0.05)
    
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["orange"]
    shape.line.color.rgb = COLORS["orange"]

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── SLIDE 01: COVER ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_background(slide, COLORS["black"])

# Top bar
add_tag(slide, "lablab.ai · April 2026", x_pos=8, y_pos=0.3)

# Main title
left = Inches(0.5)
top = Inches(1.8)
width = Inches(9)
height = Inches(2.5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "TRUST\nTHROUGH"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]
p.line_spacing = 1.2

p = text_frame.add_paragraph()
p.text = "PROOF."
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

add_divider(slide, x_pos=0.5, y_pos=4.5)

# Subtitle
left = Inches(0.5)
top = Inches(4.7)
width = Inches(8)
height = Inches(1.5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "A multi-agent AI system that earns on-chain reputation by knowing when not to trade — and proves every decision permanently."
p.font.size = Pt(14)
p.font.color.rgb = COLORS["dark"]
p.line_spacing = 1.4

# Bottom stats
stats = [
    ("$55K", "Prize Pool"),
    ("ERC-8004", "On-Chain Trust"),
    ("Kraken CLI", "CEX Execution"),
    ("LangGraph", "Agent Orchestration"),
]

start_x = 0.5
for i, (value, label) in enumerate(stats):
    x = start_x + (i * 2.35)
    
    # Value
    left = Inches(x)
    top = Inches(6.5)
    width = Inches(2)
    height = Inches(0.4)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    p = text_box.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    
    # Label
    top = Inches(6.95)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    p = text_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["dark"]

# ── SLIDE 02: PROBLEM ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS["black"])

add_tag(slide, "01 / The Problem", y_pos=0.3)

# Title
left = Inches(0.5)
top = Inches(1.2)
width = Inches(4.5)
height = Inches(2)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "DEFI BOTS ARE"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]

p = text_frame.add_paragraph()
p.text = "BLACK BOXES."
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

# Subtitle for problem
left = Inches(0.5)
top = Inches(3.5)
width = Inches(4.5)
height = Inches(1)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "Capital moves. Profits or losses happen. Nobody can verify why any decision was made — until money is already lost."
p.font.size = Pt(13)
p.font.color.rgb = COLORS["dark"]
p.line_spacing = 1.4

# Problem cards
problems = [
    ("⊘", "No Audit Trail", "Trade decisions are opaque. There's no verifiable record of agent reasoning."),
    ("⊗", "Trust Is Assumed", "Users can't distinguish competent agents from reckless ones before losses occur."),
    ("⊙", "No Reputation System", "Every deployment starts from zero. Good behavior isn't rewarded or portable."),
    ("◉", "Risk Is Hidden", "Drawdown controls are self-reported. There's no independent validation."),
]

right_col_x = 5.3
for i, (icon, title, desc) in enumerate(problems):
    y = 1.2 + (i * 1.4)
    
    # Card background
    left = Inches(right_col_x)
    top = Inches(y)
    width = Inches(4.2)
    height = Inches(1.2)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["charcoal"]
    shape.line.color.rgb = COLORS["dimRed"]
    shape.line.width = Pt(1)
    
    # Icon
    text_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.05), Inches(0.3), Inches(0.3))
    p = text_box.text_frame.paragraphs[0]
    p.text = icon
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS["orange"]
    
    # Title
    text_box = slide.shapes.add_textbox(left + Inches(0.55), top + Inches(0.05), Inches(3.4), Inches(0.25))
    p = text_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS["cream"]
    
    # Description
    text_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.35), Inches(3.9), Inches(0.75))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["dark"]
    p.line_spacing = 1.3

# ── SLIDE 03: SOLUTION ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS["charcoal"])

add_tag(slide, "02 / The Solution", y_pos=0.3)

# Left: Title
left = Inches(0.5)
top = Inches(1.2)
width = Inches(4.5)
height = Inches(2.5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "APEX EARNS TRUST BY"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]

p = text_frame.add_paragraph()
p.text = "REFUSING."
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

add_divider(slide, x_pos=0.5, y_pos=3.8)

left = Inches(0.5)
top = Inches(4.1)
width = Inches(4.5)
height = Inches(2.5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "Every trade decision — approved or refused — is a permanent, verifiable on-chain artifact via ERC-8004 registries."
p.font.size = Pt(12)
p.font.color.rgb = COLORS["dark"]
p.line_spacing = 1.4

p = text_frame.add_paragraph()
p.space_before = Pt(8)
p.text = "The Guardian agent gains reputation specifically for correct vetoes. The longer APEX runs, the more trustworthy it becomes — provably, on-chain, auditable by anyone."
p.font.size = Pt(12)
p.font.color.rgb = COLORS["dark"]
p.line_spacing = 1.4

# Right: Comparison
left = Inches(5.3)
top = Inches(1.3)
width = Inches(4.2)
height = Inches(5)

text_box = slide.shapes.add_textbox(left, top, Inches(4.2), Inches(0.5))
p = text_box.text_frame.paragraphs[0]
p.text = "The Inversion"
p.font.size = Pt(11)
p.font.color.rgb = COLORS["dark"]

# Comparison rows
comparisons = [
    ("Every Other Bot", "Trades aggressively. Trust assumed. Black box.", "#f87171", 30),
    ("APEX", "Earns reputation by refusing bad trades. Every decision on-chain.", COLORS["orange"], 90),
]

for i, (label, desc, color, pct) in enumerate(comparisons):
    y_offset = 2.0 + (i * 2.2)
    
    # Label
    text_box = slide.shapes.add_textbox(left, Inches(y_offset), Inches(4.2), Inches(0.25))
    p = text_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.bold = True
    if color == COLORS["orange"]:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = RGBColor(0xf8, 0x71, 0x71)
    
    # Description
    text_box = slide.shapes.add_textbox(left, Inches(y_offset + 0.25), Inches(4.2), Inches(0.6))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["dark"]
    p.line_spacing = 1.3
    
    # Progress bar background
    bar_shape = slide.shapes.add_shape(1, left, Inches(y_offset + 0.95), Inches(4.2), Inches(0.08))
    bar_shape.fill.solid()
    bar_shape.fill.fore_color.rgb = COLORS["dimRed"]
    bar_shape.line.color.rgb = COLORS["dimRed"]
    
    # Progress bar fill
    fill_pct = pct / 100.0
    bar_fill = slide.shapes.add_shape(1, left, Inches(y_offset + 0.95), Inches(4.2 * fill_pct), Inches(0.08))
    bar_fill.fill.solid()
    if color == COLORS["orange"]:
        bar_fill.fill.fore_color.rgb = COLORS["orange"]
    else:
        bar_fill.fill.fore_color.rgb = RGBColor(0xf8, 0x71, 0x71)
    bar_fill.line.color.rgb = bar_fill.fill.fore_color.rgb

# Key insight box
left = Inches(5.3)
top = Inches(6.5)
width = Inches(4.2)
height = Inches(0.8)
shape = slide.shapes.add_shape(1, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COLORS["black"]
shape.line.color.rgb = COLORS["dimRed"]
shape.line.width = Pt(1)

text_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.05), Inches(3.9), Inches(0.7))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "Key Insight: Risk-adjusted profitability, drawdown control, and validation quality — not raw PnL. APEX is engineered for exactly these criteria."
p.font.size = Pt(9)
p.font.color.rgb = COLORS["cream"]
p.line_spacing = 1.2

# ── SLIDE 04: HOW IT WORKS ───────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS["black"])

add_tag(slide, "03 / How It Works", y_pos=0.3)

left = Inches(0.5)
top = Inches(1.1)
width = Inches(9)
height = Inches(0.6)
text_box = slide.shapes.add_textbox(left, top, width, height)
p = text_box.text_frame.paragraphs[0]
p.text = "THE FOUR-AGENT PIPELINE"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]

# Pipeline boxes
agents = [
    ("01", "SCOUT", "Market Intel", "Scans Aave, Compound, Curve via MCP tools. Pulls volatility + sentiment.", "#60a5fa", False),
    ("02", "STRATEGIST", "Intent Gen", "Ranks opportunities. Kelly-criterion sizing. EIP-712 signed trade intents.", "#a78bfa", False),
    ("03 ⭐", "GUARDIAN", "Risk Gate", "6 veto thresholds. Temp = 0.0. Posts on-chain attestation for EVERY decision.", COLORS["orange"], True),
    ("04", "EXECUTOR", "Execution", "Real trades via Kraken CLI + Surge Risk Router. PnL recorded on-chain.", "#34d399", False),
]

pipeline_y = 2.0
for i, (num, name, role, desc, color, is_active) in enumerate(agents):
    x = 0.5 + (i * 2.4)
    
    # Box background
    box_width = Inches(2.3)
    box_height = Inches(2.8)
    left = Inches(x)
    top = Inches(pipeline_y)
    shape = slide.shapes.add_shape(1, left, top, box_width, box_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["charcoal"] if is_active else RGBColor(0, 0, 0)
    shape.line.color.rgb = COLORS["dimRed"]
    shape.line.width = Pt(1)
    
    # Number
    text_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), Inches(2.1), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x60, 0xa5, 0xfa) if color == "#60a5fa" else RGBColor(0xa7, 0x8b, 0xfa) if color == "#a78bfa" else COLORS["orange"] if is_active else RGBColor(0x34, 0xd3, 0x99)
    
    # Name
    text_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.35), Inches(2.1), Inches(0.4))
    p = text_box.text_frame.paragraphs[0]
    p.text = name
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"] if is_active else COLORS["cream"]
    
    # Role
    text_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.8), Inches(2.1), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = role
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS["dark"]
    
    # Description
    text_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(1.05), Inches(2.1), Inches(1.4))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS["dark"]
    p.line_spacing = 1.2
    
    # Status pill
    if is_active:
        pill = slide.shapes.add_shape(1, left + Inches(0.1), top + Inches(2.5), Inches(2.1), Inches(0.2))
        pill.fill.solid()
        pill.fill.fore_color.rgb = COLORS["dimRed"]
        pill.line.color.rgb = COLORS["dark"]
        pill.line.width = Pt(1)
        text_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(2.5), Inches(2.1), Inches(0.2))
        p = text_box.text_frame.paragraphs[0]
        p.text = "Earns rep for vetoes"
        p.font.size = Pt(7)
        p.font.color.rgb = COLORS["orange"]
        p.alignment = PP_ALIGN.CENTER

# ERC-8004 Layer
left = Inches(0.5)
top = Inches(5.1)
width = Inches(9)
height = Inches(0.6)
shape = slide.shapes.add_shape(1, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COLORS["charcoal"]
shape.line.color.rgb = COLORS["dimRed"]
shape.line.width = Pt(1)

text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.05), Inches(1.5), Inches(0.5))
p = text_box.text_frame.paragraphs[0]
p.text = "ERC-8004 Layer →"
p.font.size = Pt(9)
p.font.color.rgb = COLORS["orange"]

# ERC-8004 items
erc_items = [
    ("Identity Registry", "4 Agent NFTs"),
    ("Reputation Registry", "Every Decision Attested"),
    ("Validation Registry", "Vetoes Verified Post-Hoc"),
]

erc_x = 2.2
for i, (t, v) in enumerate(erc_items):
    text_box = slide.shapes.add_textbox(Inches(erc_x + i * 2.4), top + Inches(0.1), Inches(2.3), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = t
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["cream"]
    
    text_box = slide.shapes.add_textbox(Inches(erc_x + i * 2.4), top + Inches(0.35), Inches(2.3), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = v
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS["dark"]

# ── SLIDE 05: REAL TRADES ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS["charcoal"])

add_tag(slide, "04 / Real Trade Execution", y_pos=0.3)

# Left side: Title + Features
left = Inches(0.5)
top = Inches(1.1)
width = Inches(4.5)
height = Inches(2)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "REAL MONEY."
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]

p = text_frame.add_paragraph()
p.text = "REAL PROOF."
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

# Features
features = [
    ("Kraken CLI", "Direct CEX execution. Real fills. Trade history verifiable in Kraken Pro."),
    ("Surge Risk Router", "On-chain DeFi execution. Every tx on BaseScan. Provable stablecoin yield."),
    ("Dual Settlement", "Same agent cycle executes on both rails — eligible for both prize tracks."),
    ("Live PnL", "Dashboard updates within 500ms of trade confirmation. Net PnL shown."),
]

feature_y = 3.5
for i, (label, desc) in enumerate(features):
    y = feature_y + (i * 0.85)
    
    # Icon
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(0.3), Inches(0.3))
    p = text_box.text_frame.paragraphs[0]
    p.text = "◈"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS["orange"]
    
    # Label
    text_box = slide.shapes.add_textbox(Inches(0.95), Inches(y), Inches(4.05), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS["cream"]
    
    # Description
    text_box = slide.shapes.add_textbox(Inches(0.95), Inches(y + 0.22), Inches(4.05), Inches(0.55))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["dark"]
    p.line_spacing = 1.2

# Right side: Live Session Box
left = Inches(5.3)
top = Inches(1.1)
width = Inches(4.2)
height = Inches(5.9)
shape = slide.shapes.add_shape(1, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COLORS["black"]
shape.line.color.rgb = COLORS["dimRed"]
shape.line.width = Pt(1)

# Header
text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), Inches(3.8), Inches(0.25))
p = text_box.text_frame.paragraphs[0]
p.text = "Live Session · APEX Wallet"
p.font.size = Pt(9)
p.font.color.rgb = COLORS["dark"]

# PnL amount
text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5), Inches(3.8), Inches(0.6))
p = text_box.text_frame.paragraphs[0]
p.text = "+$142.56"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = COLORS["green"]
p.alignment = PP_ALIGN.CENTER

text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.1), Inches(3.8), Inches(0.2))
p = text_box.text_frame.paragraphs[0]
p.text = "Session Net PnL"
p.font.size = Pt(8)
p.font.color.rgb = COLORS["dark"]
p.alignment = PP_ALIGN.CENTER

# Stats grid
stats = [("8", "Trades Executed"), ("5", "Vetoed"), ("$2,380", "Volume"), ("0.9%", "Drawdown")]
for i, (value, label) in enumerate(stats):
    col = i % 2
    row = i // 2
    stat_x = left + Inches(0.2 + col * 1.9)
    stat_y = top + Inches(1.5 + row * 0.7)
    
    stat_box = slide.shapes.add_shape(1, stat_x, stat_y, Inches(1.8), Inches(0.6))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = COLORS["charcoal"]
    stat_box.line.color.rgb = COLORS["dimRed"]
    stat_box.line.width = Pt(1)
    
    text_box = slide.shapes.add_textbox(stat_x + Inches(0.05), stat_y + Inches(0.05), Inches(1.7), Inches(0.25))
    p = text_box.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    p.alignment = PP_ALIGN.CENTER
    
    text_box = slide.shapes.add_textbox(stat_x + Inches(0.05), stat_y + Inches(0.3), Inches(1.7), Inches(0.25))
    p = text_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(7)
    p.font.color.rgb = COLORS["dark"]
    p.alignment = PP_ALIGN.CENTER

# Trade list
trades = [
    ("USDC/USD", "BUY", "+$11.20", COLORS["green"]),
    ("BTC/USD", "VETOED", "—", COLORS["orange"]),
    ("ETH/USD", "BUY", "+$7.84", COLORS["green"]),
]

trade_y = 3.5
for i, (pair, side, pnl, color) in enumerate(trades):
    y = trade_y + (i * 0.45)
    
    text_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(y), Inches(1.5), Inches(0.3))
    p = text_box.text_frame.paragraphs[0]
    p.text = pair
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["cream"]
    
    text_box = slide.shapes.add_textbox(left + Inches(1.7), Inches(y), Inches(1), Inches(0.3))
    p = text_box.text_frame.paragraphs[0]
    p.text = side
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["dark"]
    p.alignment = PP_ALIGN.CENTER
    
    text_box = slide.shapes.add_textbox(left + Inches(2.7), Inches(y), Inches(1.3), Inches(0.3))
    p = text_box.text_frame.paragraphs[0]
    p.text = pnl
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.RIGHT

# ── SLIDE 06: SPONSOR FIT ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS["black"])

add_tag(slide, "05 / Sponsor Alignment", y_pos=0.3)

left = Inches(0.5)
top = Inches(1.1)
width = Inches(9)
height = Inches(0.8)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "BOTH CHALLENGES."
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]

p = text_frame.add_paragraph()
p.text = "ONE SUBMISSION."
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

# Sponsor boxes
sponsors = [
    {
        "name": "SURGE",
        "pool": "$50K POOL",
        "items": [
            "ERC-8004 Identity Registry — 4 agent NFTs minted",
            "ERC-8004 Reputation Registry — every decision attested",
            "ERC-8004 Validation Registry — veto correctness verified",
            "EIP-712 typed data signatures for all trade intents",
            "EIP-1271 smart-contract wallet support",
            "Surge Risk Router — whitelisted DEX execution",
            "Capital sandbox vault — stablecoin-denominated PnL",
            "Risk-adjusted ranking: drawdown < 2% target",
        ]
    },
    {
        "name": "KRAKEN",
        "pool": "$5K POOL",
        "items": [
            "Kraken CLI binary — real orders, real fills",
            "Read-only API key linked for leaderboard verification",
            "Autonomous workflow — no human intervention",
            "Market + limit order support",
            "BTC/USD, ETH/USD, SOL/USD, USDC/USD pairs",
            "Social engagement — X posts during build period",
            "Kraken Pro UI shows agent-executed trade history",
            "Combined submission: ERC-8004 identity on Kraken trades",
        ]
    },
]

for idx, sponsor in enumerate(sponsors):
    x = 0.5 + (idx * 5)
    
    # Box
    left = Inches(x)
    top = Inches(2.1)
    width = Inches(4.5)
    height = Inches(4.5)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["charcoal"]
    shape.line.color.rgb = COLORS["dimRed"]
    shape.line.width = Pt(1)
    
    # Header
    header_box = slide.shapes.add_shape(1, left, top, width, Inches(0.5))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = COLORS["black"]
    
    text_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.05), Inches(1.5), Inches(0.4))
    p = text_box.text_frame.paragraphs[0]
    p.text = sponsor["name"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    
    text_box = slide.shapes.add_textbox(left + Inches(2.5), top + Inches(0.08), Inches(1.9), Inches(0.35))
    shape_tag = slide.shapes.add_shape(1, left + Inches(2.5), top + Inches(0.08), Inches(1.9), Inches(0.35))
    shape_tag.fill.solid()
    shape_tag.fill.fore_color.rgb = COLORS["dimRed"]
    shape_tag.line.color.rgb = COLORS["dark"]
    shape_tag.line.width = Pt(1)
    
    text_box = slide.shapes.add_textbox(left + Inches(2.5), top + Inches(0.08), Inches(1.9), Inches(0.35))
    p = text_box.text_frame.paragraphs[0]
    p.text = sponsor["pool"]
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["orange"]
    p.alignment = PP_ALIGN.CENTER
    
    # Items
    items_y = top + Inches(0.65)
    for i, item in enumerate(sponsor["items"]):
        item_y = items_y + Inches(i * 0.43)
        text_box = slide.shapes.add_textbox(left + Inches(0.15), item_y, Inches(4.2), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = "✓ " + item
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS["cream"]
        p.line_spacing = 1.1

# Combined pool banner
left = Inches(0.5)
top = Inches(6.7)
width = Inches(9)
height = Inches(0.6)
shape = slide.shapes.add_shape(1, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COLORS["dimRed"]
shape.line.color.rgb = COLORS["dimRed"]

text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08), Inches(7.5), Inches(0.45))
p = text_box.text_frame.paragraphs[0]
p.text = "Combined submission = eligible for BOTH prize pools simultaneously"
p.font.size = Pt(11)
p.font.color.rgb = COLORS["cream"]

text_box = slide.shapes.add_textbox(left + Inches(7.8), top + Inches(0.05), Inches(1.2), Inches(0.5))
p = text_box.text_frame.paragraphs[0]
p.text = "$55,000"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]
p.alignment = PP_ALIGN.RIGHT

# ── SLIDE 07: TECH STACK ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS["charcoal"])

add_tag(slide, "06 / Tech Stack", y_pos=0.3)

# Left: Title
left = Inches(0.5)
top = Inches(1.1)
width = Inches(4.5)
height = Inches(2)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "PRODUCTION-GRADE"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]

p = text_frame.add_paragraph()
p.text = "STACK."
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

# Tech table
table_items = [
    ("Agents", "LangGraph 0.2", "Hierarchical state machine"),
    ("LLM", "Groq · Llama 3.3", "750 tok/s — demo-grade speed"),
    ("Trust", "ERC-8004", "On-chain identity + reputation"),
    ("DeFi Exec", "Surge Risk Router", "Whitelisted sandbox vault"),
    ("CEX Exec", "Kraken CLI", "Real orders, dual prize"),
    ("Payments", "x402 Protocol", "Agent micropayments"),
    ("Frontend", "Next.js 15", "Live PnL dashboard"),
    ("Observability", "LangSmith", "Full agent trace logs"),
]

table_y = 3.5
for i, (layer, tech, why) in enumerate(table_items):
    y = table_y + (i * 0.35)
    
    # Layer
    text_box = slide.shapes.add_textbox(left + Inches(0), Inches(y), Inches(1.4), Inches(0.3))
    p = text_box.text_frame.paragraphs[0]
    p.text = layer
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["cream"]
    p.font.bold = True
    
    # Tech
    text_box = slide.shapes.add_textbox(left + Inches(1.5), Inches(y), Inches(1.5), Inches(0.3))
    p = text_box.text_frame.paragraphs[0]
    p.text = tech
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["cream"]
    
    # Why
    text_box = slide.shapes.add_textbox(left + Inches(3.1), Inches(y), Inches(1.9), Inches(0.3))
    p = text_box.text_frame.paragraphs[0]
    p.text = why
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["dark"]

# Right: Architecture
right_x = 5.3

text_box = slide.shapes.add_textbox(Inches(right_x), Inches(1.2), Inches(4.2), Inches(0.35))
p = text_box.text_frame.paragraphs[0]
p.text = "Architecture"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS["dark"]

arch_items = [
    ("Next.js Dashboard", "Live PnL · Veto Log · Agent Registry"),
    ("FastAPI + SSE", "Real-time event stream to UI"),
    ("LangGraph Pipeline", "Scout → Strategist → Guardian → Executor"),
    ("Execution Layer", "Surge (DeFi) + Kraken CLI (CEX)"),
    ("ERC-8004 Layer", "Identity · Reputation · Validation"),
    ("IPFS + Base Sepolia", "Evidence storage + on-chain proof"),
]

arch_y = 1.8
for i, (layer, sub) in enumerate(arch_items):
    y = arch_y + (i * 0.88)
    
    # Number box
    num_box = slide.shapes.add_shape(1, Inches(right_x), Inches(y), Inches(0.35), Inches(0.35))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = COLORS["dimRed"]
    num_box.line.color.rgb = COLORS["dimRed"]
    
    text_box = slide.shapes.add_textbox(Inches(right_x), Inches(y), Inches(0.35), Inches(0.35))
    p = text_box.text_frame.paragraphs[0]
    p.text = str(i + 1).zfill(2)
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["orange"]
    p.alignment = PP_ALIGN.CENTER
    
    # Layer name
    text_box = slide.shapes.add_textbox(Inches(right_x + 0.45), Inches(y), Inches(3.7), Inches(0.22))
    p = text_box.text_frame.paragraphs[0]
    p.text = layer
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS["cream"]
    
    # Sub
    text_box = slide.shapes.add_textbox(Inches(right_x + 0.45), Inches(y + 0.25), Inches(3.7), Inches(0.35))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = sub
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS["dark"]
    p.line_spacing = 1.1

# ── SLIDE 08: BUILD PLAN ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS["black"])

add_tag(slide, "07 / Build Plan", y_pos=0.3)

# Title
left = Inches(0.5)
top = Inches(1.1)
width = Inches(4.5)
height = Inches(2)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "14 DAYS."
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]

p = text_frame.add_paragraph()
p.text = "EVERY DAY COUNTS."
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

# Timeline
timeline = [
    ("Day 1–2", "Core Execution", "Kraken CLI integration + real test trade. Surge vault connection."),
    ("Day 3", "ERC-8004 On-Chain", "4 agent NFTs minted. Validator contract deployed. Reputation signals flowing."),
    ("Day 4–5", "PnL Engine", "Real-time PnL calculation. Circuit breakers. Session manager."),
    ("Day 6–8", "Live Dashboard", "PnL chart, veto log, trade history, emergency stop UI."),
    ("Day 9–10", "Autonomous Loop", "60s cycle daemon. Health checks. 2-hour stability test."),
    ("Day 11–12", "Demo & Video", "3-minute demo showing real trade + real veto on-chain."),
    ("Day 13–14", "Submit", "README, X post, lablab.ai submission, early.surge.xyz."),
]

timeline_y = 3.2
for i, (day, title, desc) in enumerate(timeline):
    y = timeline_y + (i * 0.5)
    
    # Dot
    dot_x = Inches(0.6)
    dot_y = Inches(y + 0.15)
    dot = slide.shapes.add_shape(1, dot_x, dot_y, Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = COLORS["orange"]
    dot.line.color.rgb = COLORS["black"]
    dot.line.width = Pt(2)
    
    # Day
    text_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(1.2), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = day
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    
    # Title
    text_box = slide.shapes.add_textbox(Inches(1), Inches(y + 0.2), Inches(1.2), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["cream"]
    p.font.bold = True
    
    # Description
    text_box = slide.shapes.add_textbox(Inches(2.3), Inches(y), Inches(2.7), Inches(0.4))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS["dark"]
    p.line_spacing = 1.2

# Right: Scoring
right_x = 5.3
text_box = slide.shapes.add_textbox(Inches(right_x), Inches(1.2), Inches(4.2), Inches(0.3))
p = text_box.text_frame.paragraphs[0]
p.text = "Judging Criteria Alignment"
p.font.size = Pt(11)
p.font.color.rgb = COLORS["dark"]

scoring = [
    ("Risk-Adjusted PnL", 95, "Kelly criterion sizing, drawdown control"),
    ("Drawdown Control", 99, "Circuit breaker at 8%, Guardian at 5%"),
    ("Validation Quality", 95, "Validator contract, IPFS evidence, all 3 registries"),
    ("ERC-8004 Integration", 100, "Identity + Reputation + Validation fully used"),
    ("Kraken CLI Usage", 95, "Real orders, autonomous, leaderboard-linked"),
    ("Code Quality", 90, "TypeScript + Python, MIT license, full README"),
]

score_y = 1.7
for i, (label, score, note) in enumerate(scoring):
    y = score_y + (i * 0.82)
    
    # Label
    text_box = slide.shapes.add_textbox(Inches(right_x), Inches(y), Inches(3.5), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["cream"]
    p.font.bold = True
    
    # Score
    text_box = slide.shapes.add_textbox(Inches(right_x + 3.6), Inches(y), Inches(0.6), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = f"{score}%"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    p.alignment = PP_ALIGN.RIGHT
    
    # Bar background
    bar_bg = slide.shapes.add_shape(1, Inches(right_x), Inches(y + 0.22), Inches(4.2), Inches(0.08))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = COLORS["dimRed"]
    bar_bg.line.color.rgb = COLORS["dimRed"]
    
    # Bar fill
    fill_pct = score / 100.0
    bar_fill = slide.shapes.add_shape(1, Inches(right_x), Inches(y + 0.22), Inches(4.2 * fill_pct), Inches(0.08))
    bar_fill.fill.solid()
    bar_fill.fill.fore_color.rgb = COLORS["orange"]
    bar_fill.line.color.rgb = COLORS["orange"]
    
    # Note
    text_box = slide.shapes.add_textbox(Inches(right_x), Inches(y + 0.35), Inches(4.2), Inches(0.35))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = note
    p.font.size = Pt(7)
    p.font.color.rgb = COLORS["dark"]
    p.line_spacing = 1.1

# ── SLIDE 09: DEMO ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS["charcoal"])

add_tag(slide, "08 / The Demo Moment", y_pos=0.3)

# Title
left = Inches(0.5)
top = Inches(1.1)
width = Inches(4.5)
height = Inches(1.5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "THE 3-MIN"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]

p = text_frame.add_paragraph()
p.text = "SEQUENCE."
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

# Description
left = Inches(0.5)
top = Inches(2.8)
width = Inches(4.5)
height = Inches(0.6)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "Every second is choreographed. The veto moment is the centrepiece — a real on-chain transaction that judges can verify independently."
p.font.size = Pt(11)
p.font.color.rgb = COLORS["dark"]
p.line_spacing = 1.4

# Demo timeline
demo_frames = [
    ("0:00", "Hook", "Show 4 agents with ERC-8004 identity links on BaseScan"),
    ("0:20", "Pipeline Runs", "Scout fetches real Aave/Curve data. Strategist generates signed intent."),
    ("0:50", "VETO FIRES 🛡️", "Guardian blocks on volatility. Red banner. On-chain tx appears on BaseScan LIVE."),
    ("1:20", "Evidence Proof", "Click BaseScan tx → show IPFS evidence payload. Permanent record."),
    ("1:35", "Rep Increases", "Guardian reputation ticks up: 84.2% → 84.5%. Earns for refusing."),
    ("1:50", "Real Trade Executes", "Next cycle: Guardian approves. Kraken order fills. Surge tx confirms."),
    ("2:20", "PnL Updates", "+$11.20 appears in dashboard < 1 second after confirmation."),
    ("2:45", "Punchline", '"APEX doesn\'t just trade well. It proves why."'),
]

demo_y = 3.6
for i, (t, label, action) in enumerate(demo_frames):
    if i < 4:
        y = demo_y + (i * 0.75)
        x = 0.5
    else:
        y = demo_y + ((i - 4) * 0.75)
        x = 5.3
    
    # Time
    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.5), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = t
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    
    # Label
    text_box = slide.shapes.add_textbox(Inches(x + 0.6), Inches(y), Inches(1.2), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLORS["cream"]
    
    # Action
    text_box = slide.shapes.add_textbox(Inches(x + 1.9), Inches(y), Inches(3.4), Inches(0.7))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = action
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS["dark"]
    p.line_spacing = 1.2

# ── SLIDE 10: CTA ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS["black"])

# Main title
left = Inches(0.5)
top = Inches(1.8)
width = Inches(9)
height = Inches(2.5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "APEX\nEARNS TRUST"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = COLORS["cream"]
p.line_spacing = 1.1

p = text_frame.add_paragraph()
p.text = "ON-CHAIN."
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

add_divider(slide, x_pos=0.5, y_pos=4.3)

# Tagline
left = Inches(0.5)
top = Inches(4.6)
width = Inches(8)
height = Inches(1.2)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "Most DeFi agents try to earn trust by making money.\nAPEX earns trust by proving it knows when not to spend yours — and every proof is on-chain, permanent, and written by the agents themselves."
p.font.size = Pt(13)
p.font.color.rgb = COLORS["dark"]
p.line_spacing = 1.5

# Bottom stats
cta_stats = [
    ("$55K", "Combined Prize Pool"),
    ("2", "Prize Challenges"),
    ("100%", "On-Chain Proof"),
    ("0%", "Drawdown Target"),
]

start_x = 0.5
for i, (value, label) in enumerate(cta_stats):
    x = start_x + (i * 2.25)
    
    left = Inches(x)
    top = Inches(6)
    width = Inches(2)
    height = Inches(0.4)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    p = text_box.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    
    top = Inches(6.45)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    p = text_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["dark"]

# Footer links
links = [
    ("GitHub", "github.com/your-team/apex"),
    ("Surge Registration", "early.surge.xyz"),
    ("Contract (Base Sepolia)", "0x8004A169...432"),
]

footer_y = 7.0
for i, (label, value) in enumerate(links):
    x = 0.5 + (i * 3.1)
    
    text_box = slide.shapes.add_textbox(Inches(x), Inches(footer_y), Inches(2.9), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS["dark"]
    p.font.bold = True
    
    text_box = slide.shapes.add_textbox(Inches(x), Inches(footer_y + 0.22), Inches(2.9), Inches(0.2))
    p = text_box.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["cream"]

# Social
text_box = slide.shapes.add_textbox(Inches(7.8), Inches(footer_y), Inches(1.7), Inches(0.2))
p = text_box.text_frame.paragraphs[0]
p.text = "@lablabai · @Surgexyz_"
p.font.size = Pt(8)
p.font.color.rgb = COLORS["orange"]
p.alignment = PP_ALIGN.RIGHT

# Save presentation
output_file = "/Users/shikhar/Apex/APEX_Pitch_Deck.pptx"
prs.save(output_file)
print(f"✓ PowerPoint deck created: {output_file}")
